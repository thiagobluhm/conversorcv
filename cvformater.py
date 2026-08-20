import json
import traceback
import time
from dotenv import load_dotenv
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import RGBColor as docxRGBColor
from pathlib import Path
import base64
import re
import os
import streamlit as st
os.chdir(os.path.abspath(os.curdir))
from docx.shared import Inches, Cm
import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as pptxRGBColor
from datetime import date
import json, re, unicodedata


class cvFormatter():
    def __init__(self):
        pass

    def log_etapa(self, mensagem):
        """Imprime uma linha de log de acompanhamento no terminal/console.

        Uso interno apenas (controle de processamento) — não aparece em nenhum
        momento na tela para o usuário do app.
        """
        agora = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        print(f"[LOG {agora}] {mensagem}")

    def validate_json(self, dados, estrutura_padrao):
        """Valida e completa o JSON com estrutura padrão."""
        for chave in estrutura_padrao:
            if chave not in dados:
                dados[chave] = estrutura_padrao[chave]
        return dados

    def create_docx_curriculo(self, arquivo_json, arquivo_saida='curriculo.docx', logo_path='Logo2.png'):
        """Cria um documento Word formatado a partir de dados de um currículo em JSON e adiciona um logo."""
        inicio = time.time()
        self.log_etapa(f"Geração do DOCX - iniciada ({arquivo_saida})")
        try:
            with open(arquivo_json, 'r', encoding='utf-8') as f:
                dados = json.load(f)

            estrutura_padrao = {
                "informacoes_pessoais": {"nome": "", "cidade": ""},
                "resumo_qualificacoes": [],
                "perfil_profissional": "Sem informações",
                "perfil_comportamental": "Sem informações",
                "vaga": "Sem informações",
                "modalidade": "Sem informações",
                "experiencia_profissional": [],
                "educacao": [],
                "certificacoes": []
            }
            dados = self.validate_json(dados, estrutura_padrao)

            doc = Document()
            estilo = doc.styles['Normal']
            estilo.font.name = 'Calibri'
            estilo.font.size = Pt(11)
            estilo.font.color.rgb = docxRGBColor(0, 0, 0)

            def adicionar_espaco():
                """Adiciona um parágrafo vazio para espaçamento."""
                doc.add_paragraph().paragraph_format.space_after = Pt(12)

            if logo_path:
                section = doc.sections[0]
                section.header_distance = Cm(0.6)

                header = section.header
                header_paragraph = header.paragraphs[0]
                run = header_paragraph.add_run()
                run.add_picture(logo_path, width=Inches(0.8))  # Ajusta o tamanho do logo
                header_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT  # Alinha à direita

            # Informações pessoais
            informacoes_pessoais = dados.get('informacoes_pessoais', {})
            nome = informacoes_pessoais.get('nome') or 'Sem informações'
            paragrafo_nome = doc.add_paragraph(nome)
            if paragrafo_nome.runs:
                nome_run = paragrafo_nome.runs[0]
                nome_run.bold = True
                nome_run.font.size = Pt(16)
            paragrafo_nome.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

            adicionar_espaco()
            contato = f"Cidade: {informacoes_pessoais.get('cidade') or 'Sem informações'}\nVaga: {dados.get('vaga') or 'Sem informações'}\nModalidade: {dados.get('modalidade') or 'Sem informações'}"
            doc.add_paragraph(contato)

            adicionar_espaco()

            # Educação
            doc.add_heading('Educação', level=2)
            lista_educacao = dados.get('educacao') or []
            if lista_educacao:
                for educacao in lista_educacao:
                    instituicao = educacao.get('instituicao') or 'Sem informações'
                    curso = educacao.get('curso') or 'Sem informações'
                    periodo = educacao.get('periodo') or 'Sem informações'

                    doc.add_paragraph(f"{instituicao}", style='Heading 3')
                    doc.add_paragraph(f"{curso} - {periodo}", style='Normal')
            else:
                doc.add_paragraph("Sem informações", style='Normal')

            adicionar_espaco()

            # Certificações
            doc.add_heading('Certificações', level=2)
            lista_certificacoes = dados.get('certificacoes') or []
            if lista_certificacoes:
                for certificacao in lista_certificacoes:
                    doc.add_paragraph(f"{certificacao}", style='List Bullet')
            else:
                doc.add_paragraph("Sem informações", style='Normal')

            adicionar_espaco()

            # ===== INÍCIO: seção Resumo de Qualificações =====
            # Removida (ainda em validação) a pedido do Tiago: as qualificações
            # técnicas agora devem aparecer embutidas no texto do Resumo
            # Profissional, em vez de numa seção própria. Código mantido
            # comentado para reativar facilmente se precisarmos voltar atrás.
            #
            # # Resumo de qualificações
            # doc.add_heading('Resumo de Qualificações', level=2)
            # lista_qualificacoes = dados.get('resumo_qualificacoes') or []
            # if lista_qualificacoes:
            #     for qualificacao in lista_qualificacoes:
            #         doc.add_paragraph(f"- {qualificacao}")
            # else:
            #     doc.add_paragraph("Sem informações", style='Normal')
            #
            # adicionar_espaco()
            # ===== FIM: seção Resumo de Qualificações =====

            # Perfil Profissional (título exibido como "Resumo Profissional")
            doc.add_heading('Resumo Profissional', level=2)
            paragrafo_resumo_profissional = doc.add_paragraph(dados.get('perfil_profissional') or "Sem informações")
            paragrafo_resumo_profissional.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY

            adicionar_espaco()

            # Perfil Comportamental
            doc.add_heading('Perfil Comportamental', level=2)
            paragrafo_perfil_comportamental = doc.add_paragraph(dados.get('perfil_comportamental') or "Sem informações")
            paragrafo_perfil_comportamental.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY

            adicionar_espaco()

            # Experiência profissional
            doc.add_heading('Experiência Profissional', level=2)
            lista_experiencias = dados.get('experiencia_profissional') or []
            if lista_experiencias:
                for experiencia in lista_experiencias:
                    empresa = experiencia.get('empresa') or 'Sem informações'
                    cargo = experiencia.get('cargo') or 'Sem informações'
                    periodo = experiencia.get('periodo') or 'Sem informações'
                    local = experiencia.get('local') or 'Sem informações'
                    atividades = experiencia.get('atividades_exercidas') or []

                    doc.add_paragraph(f"{empresa} ({local})", style='Heading 3')
                    doc.add_paragraph(f"{cargo} - {periodo}", style='Normal')

                    doc.add_paragraph("Atividades exercidas:", style='Normal')
                    if atividades:
                        for atividade in atividades:
                            doc.add_paragraph(f"{atividade}", style='List Bullet')
                    else:
                        doc.add_paragraph("Sem informações", style='List Bullet')

                    ferramentas = experiencia.get('ferramentas') or []
                    doc.add_paragraph("Ferramentas utilizadas:", style='Normal')
                    if ferramentas:
                        for ferramenta in ferramentas:
                            doc.add_paragraph(f"{ferramenta}", style='List Bullet')
                    else:
                        doc.add_paragraph("Sem informações", style='List Bullet')
            else:
                doc.add_paragraph("Sem informações", style='Normal')

            # Salvar o documento Word
            doc.save(arquivo_saida)
            duracao = time.time() - inicio
            self.log_etapa(f"Geração do DOCX - sucesso ({arquivo_saida}, {duracao:.2f}s)")
            print(f"Currículo salvo em {arquivo_saida}")

        except Exception as e:
            duracao = time.time() - inicio
            self.log_etapa(f"Geração do DOCX - ERRO ({duracao:.2f}s): {e}")
            print(f"Erro ao criar documento Word: {e}")
            print(traceback.format_exc())


    def create_parecer_pptx(
        self,
        arquivo_json: str,
        arquivo_saida: str,
        template_path: str | None = None,
        responsavel: str | None = None,
        
    ):
        """
        Preenche um PPTX a partir de um dicionário/JSON, preservando a formatação do template.
        Estratégia:
        1) Se existir shape com NOME igual à chave (ex.: 'nome', 'modalidade'), escreve nela.
        2) Caso contrário, substitui placeholder {{chave}} no texto de shapes/células.
        Não altera estilos (tamanhos, bullets, espaçamentos): tudo vem do template.
        """

        # ---------- utilitários ----------
        import re, json, unicodedata
        from pathlib import Path
        from datetime import date
        from pptx import Presentation

        def _normalize(s: str) -> str:
            if s is None: return ""
            s = unicodedata.normalize("NFD", s)
            s = "".join(ch for ch in s if unicodedata.category(ch) != "Mn")  # remove acentos
            s = re.sub(r"\s+", " ", s.strip().lower())
            return s

        def _iter_shapes(container):
            # percorre shapes e subshapes (grupos)
            for shp in getattr(container, "shapes", []):
                yield shp
                if hasattr(shp, "shapes"):   # grupo
                    for s in _iter_shapes(shp):
                        yield s

        def _clone_paragraph_style(src_paragraph, dst_paragraph):
            """
            Clona o pPr (Paragraph Properties) do parágrafo fonte para o destino,
            garantindo espaçamento/bullets/entrelinhas idênticos.
            """
            try:
                from copy import deepcopy
                src_pPr = getattr(src_paragraph._p, "pPr", None)
                if src_pPr is not None:
                    # garante que o destino tenha pPr e substitui pelo do fonte
                    dst_pPr = dst_paragraph._p.get_or_add_pPr()
                    dst_paragraph._p.remove(dst_pPr)
                    dst_paragraph._p.append(deepcopy(src_pPr))
            except Exception:
                # se a API/estruturas variarem, seguimos com o padrão herdado
                pass

        def _set_text_preservando_estilo(tf, text: str):
            """
            Preserva exatamente o estilo do template:
            - Reutiliza o 1º run do 1º parágrafo (mantém rPr: fonte/tamanho/estilo).
            - Para múltiplas linhas, clona o parágrafo inteiro (p) do p0 para cada linha,
            mantendo pPr e rPr. Assim, 2º parágrafo fica 100% igual ao 1º.
            """
            from copy import deepcopy

            # Garante ao menos 1 parágrafo
            if not tf.paragraphs:
                p0 = tf.add_paragraph()
            else:
                p0 = tf.paragraphs[0]

            # Normaliza linhas (um parágrafo por linha)
            linhas = (text or "").split("\n")
            if not linhas:
                linhas = [""]

            # --- prepara p0 com a 1ª linha, SEM destruir seus runs ---
            # se não existir run, cria um vazio para “pegar” a rPr do parágrafo
            if not p0.runs:
                r = p0.add_run()
                r.text = ""  # recebe rPr default do parágrafo/shape

            # mantém somente o 1º run de p0 (preserva rPr dele)
            first_run = p0.runs[0]
            for r in list(p0.runs)[1:]:
                try:
                    r._r.getparent().remove(r._r)
                except Exception:
                    pass
            first_run.text = linhas[0]

            # --- remove parágrafos excedentes (se já existirem) ---
            while len(tf.paragraphs) > 1:
                p_last = tf.paragraphs[-1]
                tf._element.remove(p_last._p)

            # --- gera parágrafos restantes clonando p0 integralmente ---
            for linha in linhas[1:]:
                # clona o parágrafo inteiro (p0._p) => preserva pPr e rPr
                clone_p = deepcopy(p0._p)
                tf._element.append(clone_p)
                # obtém o objeto paragraph recém-adicionado
                p = tf.paragraphs[-1]

                # garante um run e substitui somente o texto do 1º run
                if not p.runs:
                    rr = p.add_run(); rr.text = ""
                # remove runs além do primeiro (mantendo rPr do primeiro)
                for r in list(p.runs)[1:]:
                    try:
                        r._r.getparent().remove(r._r)
                    except Exception:
                        pass
                p.runs[0].text = linha


        def _replace_placeholders_textlike(s: str, token_map, counts: dict) -> str:
            out = s or ""
            for k, v in token_map.items():
                # {{  chave  }} tolerante a espaços e case-insensitive
                pat = re.compile(r"\{\{\s*" + re.escape(k) + r"\s*\}\}", re.IGNORECASE)
                new_out, n = pat.subn(v, out)
                if n:
                    counts[k] = counts.get(k, 0) + n
                out = new_out
            return out

                # ---------- template ----------
        if template_path is None:
            template_path = str(Path(__file__).with_name("PARECER_MODELO2.pptx"))
        if not Path(template_path).exists():
            raise FileNotFoundError(f"Template PPTX não encontrado em: {template_path}")

        # ---------- carrega JSON ----------
        with open(arquivo_json, "r", encoding="utf-8") as f:
            dados = json.load(f) or {}
        d = {(k.lower() if isinstance(k, str) else k): v for k, v in dados.items()}

        # ---------- formatadores (simples; template manda no estilo) ----------
        def _as_str(v): 
            return "" if v is None else str(v)

        def _as_pipe(v):
            if isinstance(v, list):
                return " | ".join(str(x) for x in v if str(x).strip())
            return _as_str(v)

        def _as_paragraphs(v):
            # lista -> um parágrafo por item (sem linha em branco extra)
            if isinstance(v, list):
                return "\n".join(str(x) for x in v if str(x).strip())
            return _as_str(v)

        def _as_formacao(v):
            # linhas simples; o template decide bullets/estilo
            linhas = []
            for f in (v or []):
                linhas.append(
                    f"{f.get('grau','')} em {f.get('curso','')} — "
                    f"{f.get('instituicao','')} ({f.get('conclusao','')})"
                )
            return "\n".join([ln for ln in linhas if ln.strip()])

        # ---------- responsável + data (prioriza argumento da função) ----------
        resp_arg  = (responsavel or "").strip()
        resp_json = _as_str(d.get("responsavel")).strip()
        resp_final = resp_arg or resp_json or "Responsável"

        data_parecer_str = date.today().strftime("%d/%m/%Y")

        # ---------- valores a preencher ----------
        values = {
            "nome": _as_str(d.get("nome")),
            "disponibilidade": _as_str(d.get("disponibilidade")),
            "modalidade": _as_str(d.get("modalidade")),
            "dados_pessoais": _as_str(d.get("dados_pessoais")),
            "perfil_profissional": _as_paragraphs(d.get("perfil_profissional")),
            "perfil_comportamental": _as_paragraphs(d.get("perfil_comportamental")),
            "competencias": _as_pipe(d.get("competencias")),
            "formacao": _as_formacao(d.get("formacao")),

            # campos individuais
            "responsavel": resp_final,
            "data_parecer": data_parecer_str,

            # linha combinada (use shape 'responsavel_data' ou placeholder correspondente)
            "responsavel_data": f"Responsável: {resp_final}  Data do parecer: {data_parecer_str}",
        }


        # ---------- abre modelo ----------
        prs = Presentation(template_path)

        # Índice de shapes por NOME normalizado (slides + layouts + masters)
        shapes_by_name = {}
        def _indexar(container):
            for shp in _iter_shapes(container):
                nm = _normalize(getattr(shp, "name", ""))
                if nm:
                    shapes_by_name.setdefault(nm, []).append(shp)

        for slide in prs.slides: _indexar(slide)
        for master in prs.slide_masters:
            _indexar(master)
            for layout in master.slide_layouts: _indexar(layout)
        for layout in prs.slide_layouts: _indexar(layout)

        # ---------- 1ª passada: preencher por NOME de shape ----------
        filled_by_name = set()
        for key, val in values.items():
            nm = _normalize(key)
            if nm in shapes_by_name:
                for shp in shapes_by_name[nm]:
                    if getattr(shp, "has_text_frame", False):
                        _set_text_preservando_estilo(shp.text_frame, val)
                    elif getattr(shp, "has_table", False):
                        try:
                            cell = shp.table.cell(0,0)
                            _set_text_preservando_estilo(cell.text_frame, val)
                        except Exception:
                            pass
                filled_by_name.add(key)

        # ---------- 2ª passada: substituir placeholders {{chave}} onde sobrou ----------
        remaining = {k: v for k, v in values.items() if k not in filled_by_name}
        placeholder_hits = {}
        if remaining:
            # slides
            for slide in prs.slides:
                for shp in _iter_shapes(slide):
                    if getattr(shp, "has_text_frame", False):
                        updated = _replace_placeholders_textlike(shp.text_frame.text, remaining, placeholder_hits)
                        if updated != shp.text_frame.text:
                            _set_text_preservando_estilo(shp.text_frame, updated)
                    if getattr(shp, "has_table", False):
                        for row in shp.table.rows:
                            for cell in row.cells:
                                upd = _replace_placeholders_textlike(cell.text, remaining, placeholder_hits)
                                if upd != cell.text:
                                    _set_text_preservando_estilo(cell.text_frame, upd)
            # masters e layouts
            for master in prs.slide_masters:
                for shp in _iter_shapes(master):
                    if getattr(shp, "has_text_frame", False):
                        updated = _replace_placeholders_textlike(shp.text_frame.text, remaining, placeholder_hits)
                        if updated != shp.text_frame.text:
                            _set_text_preservando_estilo(shp.text_frame, updated)
                    if getattr(shp, "has_table", False):
                        for row in shp.table.rows:
                            for cell in row.cells:
                                upd = _replace_placeholders_textlike(cell.text, remaining, placeholder_hits)
                                if upd != cell.text:
                                    _set_text_preservando_estilo(cell.text_frame, upd)
            for layout in prs.slide_layouts:
                for shp in _iter_shapes(layout):
                    if getattr(shp, "has_text_frame", False):
                        updated = _replace_placeholders_textlike(shp.text_frame.text, remaining, placeholder_hits)
                        if updated != shp.text_frame.text:
                            _set_text_preservando_estilo(shp.text_frame, updated)
                    if getattr(shp, "has_table", False):
                        for row in shp.table.rows:
                            for cell in row.cells:
                                upd = _replace_placeholders_textlike(cell.text, remaining, placeholder_hits)
                                if upd != cell.text:
                                    _set_text_preservando_estilo(cell.text_frame, upd)

        # ---------- salva (sem tocar no modelo) ----------
        Path(arquivo_saida).parent.mkdir(parents=True, exist_ok=True)
        prs.save(arquivo_saida)

        # ---------- logs (opcional) ----------
        try:
            missing = [k for k in values.keys() if k not in filled_by_name and k not in placeholder_hits]
            print(f"[PPTX] Por nome: {sorted(list(filled_by_name))}")
            print(f"[PPTX] Por placeholder: {sorted([k for k in placeholder_hits.keys()])}")
            if missing:
                print("[PPTX] Sem destino no template:", ", ".join(missing))
        except Exception:
            pass

        return arquivo_saida


    def melhorar_texto(self, texto):
        """Envia um texto para a IA melhorar gramática, clareza e coesão, mantendo as mesmas informações."""
        load_dotenv()
        chave_api = os.getenv('OPENAI_API_KEY')

        if not chave_api:
            st.error("Chave da API OpenAI não encontrada.")
            return texto

        endpoint = "https://fdrybluhm.services.ai.azure.com/openai/v1"
        deployment_name = "gpt-5.6-luna"
        client = OpenAI(base_url=endpoint, api_key=chave_api)

        prompt_melhoria = f"""Melhore a gramática, clareza e coesão do texto abaixo, mantendo exatamente as mesmas informações e o
                            mesmo sentido original. Não invente, não remova e não resuma conteúdo — apenas reescreva de forma mais
                            bem escrita. Use linguagem formal, porém de fácil entendimento. As informações são referentes a candidatos
                            a oportunidades de trabalho no Grupo Portfolio. Retorne APENAS o texto melhorado, sem aspas, sem comentários
                            e sem explicações adicionais.

TEXTO:
{texto}
"""

        try:
            response = client.responses.create(
                model=deployment_name,
                input=[
                    {"role": "system", "content": "Você é um especialista em revisão de textos em português do Brasil, focado em gramática, clareza e coesão, e trabalha com processos seletivos do Grupo Portfolio."},
                    {"role": "user", "content": prompt_melhoria}
                ],
                reasoning={"effort": "low"},
                max_output_tokens=2048
            )

            texto_melhorado = response.output_text.strip()

            usage = getattr(response, "usage", None)
            if usage:
                print(f"[TOKENS] melhorar_texto - entrada={usage.input_tokens} saida={usage.output_tokens} total={usage.total_tokens}")

            return texto_melhorado or texto
        except Exception as e:
            print(f"Erro ao melhorar texto com a API OpenAI: {e}")
            st.error("Não foi possível melhorar o texto agora. Tente novamente.")
            return texto

    def process_text_curriculo(self, texto):
        """Processa o texto e retorna JSON estruturado."""
        inicio = time.time()
        self.log_etapa("Processamento do currículo (chamada à IA) - iniciado")

        load_dotenv()
        chave_api = os.getenv('OPENAI_API_KEY')

        if not chave_api:
            self.log_etapa(f"Processamento do currículo - ERRO ({time.time() - inicio:.2f}s): chave da API OpenAI não encontrada")
            st.error("Chave da API OpenAI não encontrada.")
            return {}

        endpoint = "https://fdrybluhm.services.ai.azure.com/openai/v1"
        deployment_name = "gpt-5.6-luna"
        client = OpenAI(base_url=endpoint, api_key=chave_api)

        modelo_prompt = f"""
                            TEXTO DO CURRÍCULO:
                            {texto}

                            Campos esperados e explicações:
                            1. **informacoes_pessoais**: 
                                Contém as informações pessoais do candidato, incluindo:
                                - "nome": Nome completo do candidato.
                                - "cidade": Cidade e estado de residência.

                            2. **resumo_qualificacoes**:
                                Lista com as principais habilidades, competências ou realizações do candidato, como:
                                - Conhecimentos técnicos (ex.: Power BI, Python, SQL).
                                - Soft skills (ex.: liderança, trabalho em equipe).
                                - Principais realizações (ex.: "Aumentou a eficiência em X% ao implementar [projeto]").

                            3. **experiencia_profissional**:
                                Lista de experiências profissionais relevantes. Cada entrada deve conter:
                                - "empresa": Nome da empresa.
                                - "cargo": Cargo exercido.
                                - "periodo": Período de atuação (ex.: Janeiro de 2020 - Dezembro de 2022).
                                - "local": Local onde o trabalho foi realizado (ex.: Remoto ou Cidade/Estado).
                                - "atividades_exercidas": Lista de atividades e responsabilidades no cargo. Detalhe as principais contribuições e tarefas realizadas.
                                - "ferramentas": Lista das ferramentas, softwares ou tecnologias utilizadas no cargo (ex.: Power BI, Python, SQL, Tableau).

                            4. **educacao**:
                                Lista de formações acadêmicas do candidato. Cada entrada deve conter:
                                - "instituicao": Nome da instituição de ensino.
                                - "curso": Curso ou programa concluído.
                                - "periodo": Período de realização (ex.: Janeiro de 2016 - Dezembro de 2020).

                            5. **certificacoes**:
                                Lista de certificações relevantes obtidas pelo candidato. Cada entrada deve conter:
                                - Nome da certificação (ex.: "Microsoft Certified: Data Analyst Associate").
                                - Instituição que emitiu a certificação (ex.: Microsoft, AWS, etc.).

                            6. **perfil_profissional**:
                                Resumo profissional em texto corrido (sem divisão por tópicos), 1 único parágrafo com
                                aproximadamente 8 linhas, focando nas experiências profissionais, ordenado da mais
                                recente para a mais antiga. Evite redundâncias ao máximo, foque em não ser repetitivo.
                                Evite datas neste trecho. Sem adjetivações, sem verbos no imperativo, sem exorbitância
                                ou engrandecimento — seja factual e objetivo, mas venda bem o(a) candidato(a). Inclua
                                projetos, cases, resultados e entregas em destaque, incluindo o resultado gerado nas
                                empresas quando essa informação existir no currículo. Não crie nem invente números que
                                não estejam no texto original. Não inclua informações de formação acadêmica ou
                                certificações — elas já aparecem em outra seção do currículo. Estruture o parágrafo
                                assim: (1) visão geral da trajetória/anos de experiência e área de atuação; (2) cargo
                                e atuação mais recente, com contexto e responsabilidades; (3) passagens anteriores
                                relevantes, com projetos e ferramentas quando houver; (4) uma frase final conectando o
                                perfil do candidato ao tipo de vaga ou área em que ele mais se encaixa, com base na
                                própria trajetória — use um fechamento no estilo "apresenta forte aderência à vaga
                                por..." ou "seu perfil apresenta boa aderência para posições de...".

                                Exemplos de tom e estrutura esperados (use como referência de estilo — não copie o
                                conteúdo, cada currículo tem sua própria história):

                                Exemplo A: "Profissional com mais de 15 anos de experiência em tecnologia, com
                                trajetória que evoluiu de desenvolvimento de sistemas para gestão de projetos e
                                produtos digitais. Atuou recentemente como Project Manager e Product Manager na
                                Onnitech, liderando projetos de banco digital, aplicativos, integrações e canais de
                                atendimento, incluindo WhatsApp. Possui sólida experiência em gestão de escopo,
                                cronograma, riscos, indicadores e alinhamento entre áreas de negócio, tecnologia, UX e
                                parceiros. Na M. Dias Branco, gerenciou até seis projetos simultaneamente, além de
                                atuar com qualidade de software, análise de negócios e homologações. Também acumula
                                experiência em levantamento e documentação de requisitos pela Capgemini. Apresenta
                                forte aderência à vaga pela vivência em projetos digitais, gestão de stakeholders e
                                aplicação de metodologias ágeis e preditivas."

                                Exemplo B: "Alessa Carvalho possui experiência em gestão de projetos, transformação
                                digital e modernização de processos, com atuação em órgãos públicos de grande porte.
                                Atualmente é Gerente de Projetos no Ministério Público do Estado do Ceará, liderando
                                iniciativas de transformação digital, automação de processos e acompanhamento de
                                indicadores estratégicos. Anteriormente atuou como Gerente de Projetos de Produtos
                                Digitais na SEPLAG, conduzindo projetos de automação, produtos digitais e soluções
                                voltadas ao cidadão, incluindo Ceará App, Ceará Digital e Acesso Cidadão. Possui
                                experiência em metodologias ágeis, gestão de equipes multidisciplinares, stakeholders,
                                indicadores, riscos, contratos e recursos de projetos. Também atuou como Analista de
                                Sistemas e Analista de Processos e Projetos, com foco em governança, implantação de
                                sistemas e melhoria de processos. Seu perfil apresenta boa aderência para posições de
                                Gerente de Projetos de TI, especialmente em ambientes de transformação digital e
                                gestão estratégica."

                                Exemplo C: "Profissional com mais de 10 anos de experiência em Análise de Negócios,
                                atuando na interface entre áreas de negócio e tecnologia nos segmentos de Telecom,
                                Saúde e Agro. Possui sólida experiência em levantamento e documentação de requisitos,
                                gestão e priorização de backlog, elaboração de User Stories e acompanhamento de
                                projetos em ambientes ágeis. Atuou como Analista de Negócios e Product Owner,
                                conduzindo projetos de ponta a ponta, desde a descoberta das necessidades até a
                                homologação e entrega das soluções. Destaca-se pela participação em iniciativas de
                                transformação digital, projetos centrados no usuário e desenvolvimento de aplicações
                                corporativas, incluindo soluções para a área da saúde. Apresenta forte aderência à
                                vaga por sua experiência em relacionamento com stakeholders, mapeamento de demandas,
                                validação de requisitos e apoio à implementação de soluções tecnológicas alinhadas
                                aos objetivos do negócio."

                            Formato esperado do JSON de saída:
                            {{
                                "informacoes_pessoais": {{
                                    "nome": "",
                                    "cidade": "",
                                }},
                                "resumo_qualificacoes": [
                                    "Resumo 1",
                                    "Resumo 2"
                                ],
                                "experiencia_profissional": [
                                    {{
                                        "empresa": "Empresa X",
                                        "cargo": "Cargo Y",
                                        "periodo": "Janeiro de 2020 - Dezembro de 2022",
                                        "local": "Cidade/Estado",
                                        "atividades_exercidas": [
                                            "Atividade 1",
                                            "Atividade 2"
                                        ],
                                        "ferramentas": [
                                            "Ferramenta 1",
                                            "Ferramenta 2"
                                        ]
                                    }}
                                ],
                                "educacao": [
                                    {{
                                        "instituicao": "Instituição A",
                                        "curso": "Curso B",
                                        "periodo": "Janeiro de 2016 - Dezembro de 2020"
                                    }}
                                ],
                                "certificacoes": [
                                    "Certificação 1",
                                    "Certificação 2"
                                ],
                                "perfil_profissional": "Texto corrido do resumo profissional, em um único parágrafo."
                            }}
                            """

        try:
            response = client.responses.create(
                model=deployment_name,
                input=[
                    {"role": "system", "content": """Você é um especialista em análise de currículos e extração de informações.
                                                     Colete todas as informações possíveis, não deixe nada passar.
                                                     Dê sua resposta APENAS com o json solicitado e nada mais. NÃO ESCREVA ```json na resposta!
                    """},
                    {"role": "user", "content": modelo_prompt}
                ],
                reasoning={"effort": "low"},
                max_output_tokens=4096
            )

            conteudo = response.output_text.replace("```json", "").strip()
            # st.write(f"CONTEUDO: {conteudo}")

            usage = getattr(response, "usage", None)
            if usage:
                print(f"[TOKENS] process_text_curriculo - entrada={usage.input_tokens} saida={usage.output_tokens} total={usage.total_tokens}")

            try:
                resultado = json.loads(conteudo)
                duracao = time.time() - inicio
                nome_candidato = (resultado.get('informacoes_pessoais') or {}).get('nome', 'Sem informações')
                self.log_etapa(f"Processamento do currículo - sucesso (candidato: {nome_candidato}, {duracao:.2f}s)")
                return resultado
            except json.JSONDecodeError:
                duracao = time.time() - inicio
                self.log_etapa(f"Processamento do currículo - ERRO ({duracao:.2f}s): resposta da IA não é um JSON válido")
                print("Erro ao converter resposta da API para JSON.")
                return {}
        except Exception as e:
            duracao = time.time() - inicio
            self.log_etapa(f"Processamento do currículo - ERRO ({duracao:.2f}s): {e}")
            print(f"Erro ao processar texto com a API OpenAI: {e}")
            return {}

    def process_text_parecer(self, texto):
        """Processa o texto e retorna JSON estruturado."""
        load_dotenv()
        chave_api = os.getenv('OPENAI_API_KEY')

        if not chave_api:
            st.error("Chave da API OpenAI não encontrada.")
            return {}

        endpoint = "https://fdrybluhm.services.ai.azure.com/openai/v1"
        deployment_name = "gpt-5.6-luna"
        client = OpenAI(base_url=endpoint, api_key=chave_api)

        modelo_prompt_parecer = f"""
            TEXTO DO CURRÍCULO ORIGINAL:
            {texto}

            ### INSTRUÇÕES
            - Extraia só o que estiver presente no currículo; não invente dados.
            - Preencha **todos** os campos abaixo sempre que encontrar a informação.
            - **Formato de saída**: JSON **sem** crases, sem ```json, sem comentários.

            ### CAMPOS E PADRÕES ESPERADOS

            1. Nome
            • Nome do candidato 

            2. formacao (lista de objetos)
            • grau        → "Tecnólogo", "Bacharel", "MBA", etc.
            • curso       → Nome do curso
            • instituicao → Onde cursou
            • conclusao   → "2018", "cursando", etc.

            3. Competencias (lista de competencias)
                                Lista com as principais (no máximo 5)competências do candidato, como:
                                - Competências(ex.: Power BI, Python, SQL).

            4. perfil_profissional (listagem de 2 parágrafos, nesta ordem)
            • Parágrafo 1 – trajetória (empresas, cargos, período, volume de entregas).  
            • Parágrafo 2 – competências + projetos relevantes iniciados por verbo no infinitivo/gerúndio.

            ### EXEMPLO DE SAÍDA ESPERADA
            "Nome": "João da Silva",
            "formacao": [
                {{
                "grau": "Tecnólogo"}}]
            {{
            "formacao": [
                {{
                "grau": "Tecnólogo",
                "curso": "Análise de Sistemas",
                "instituicao": "Faculdade X",
                "conclusao": "2012"
                }},
                {{
                "grau": "MBA",
                "curso": "Gestão de Projetos",
                "instituicao": "Universidade Y",
                "conclusao": "2019"
                }}
            ],
            "competencias": [
                "Conhecimentos técnicos: Power BI, Python, SQL.",
                "Soft skills: liderança, trabalho em equipe."
            ],
            "perfil_profissional": [
                "Camila atua desde fevereiro de 2021 na empresa Sankhya como Consultora de Implantação de ERP Sênior – módulo HCM, participando de 15 projetos de implantação e conduzindo treinamentos para clientes em vários estados. Antes disso, trabalhou na Solar Coca-Cola, YDUQS e Adtalem com foco em SAP HCM, somando experiência prévia de seis anos em rotinas de departamento pessoal.",
                "Domina metodologias ágeis e Waterfall, conduz migrações de dados de sistemas legados, parametriza folha, ponto e avaliação de desempenho e implanta soluções de ERP. Implantou dois novos Centros de Distribuição e uma loja, integrou plataformas Totvs e Fortes e automatizou rotinas de importação de pedidos, entregando ganhos de produtividade em até 5 meses."
            ]
            }}
            """
        try:
            response = client.responses.create(
                model=deployment_name,
                input=[
                    {"role": "system", "content": """Você é um especialista em análise de currículos e extração de informações.
                                                     Colete todas as informações possíveis, não deixe nada passar.
                                                     Dê sua resposta APENAS com o json solicitado e nada mais. NÃO ESCREVA ```json na resposta!
                    """},
                    {"role": "user", "content": modelo_prompt_parecer}
                ],
                reasoning={"effort": "low"},
                max_output_tokens=4096
            )

            conteudo = response.output_text.replace("```json", "").strip()
            print(f"CONTEUDO: {conteudo}")

            usage = getattr(response, "usage", None)
            if usage:
                print(f"[TOKENS] process_text_parecer - entrada={usage.input_tokens} saida={usage.output_tokens} total={usage.total_tokens}")

            try:
                return json.loads(conteudo)
            except json.JSONDecodeError:
                print("Erro ao converter resposta da API para JSON.")
                return {}
        except Exception as e:
            print(f"Erro ao processar texto com a API OpenAI: {e}")
            return {}



    def extract_text_from_pdf(self, caminho_pdf):
        """Extrai o texto de um arquivo PDF."""
        inicio = time.time()
        self.log_etapa(f"Extração de texto do PDF - iniciada ({caminho_pdf})")
        try:
            leitor = PdfReader(caminho_pdf)
            texto = "".join(pagina.extract_text() for pagina in leitor.pages)
            texto = self.clear_text(texto)
            duracao = time.time() - inicio
            self.log_etapa(f"Extração de texto do PDF - sucesso ({len(texto)} caracteres, {duracao:.2f}s)")
            return texto
        except Exception as e:
            duracao = time.time() - inicio
            self.log_etapa(f"Extração de texto do PDF - ERRO ({duracao:.2f}s): {e}")
            print(f"Erro ao extrair texto do PDF: {e}")
            return ""

    def clear_text(self, texto):
        """Limpa e normaliza o texto extraído."""
        texto = re.sub(r'\s+', ' ', texto)
        texto = re.sub(r'\n*Página \d+ de \d+\n*', '', texto)
        return texto.strip()

    # Função para adicionar uma imagem de fundo a partir de um arquivo local
    def add_bg_from_local(self, image_file):
        with Path(image_file).open("rb") as file:
            encoded_string = base64.b64encode(file.read()).decode()
        st.markdown(
            f"""
            <style>
            .stApp {{
                background-color: rgba(247,247,247,0.75);
                background-size: contain;
                background-position: center;
                background-repeat: no-repeat;
                border-color: rgba(31,216,135,1) ;
    
            }}
            </style>
            """,
            unsafe_allow_html=True
        )


    def add_logo_from_local(self, logo_file):
        with Path(logo_file).open("rb") as file:
            encoded_string = base64.b64encode(file.read()).decode()

        st.markdown(
            f"""
            <style>
            [data-testid="stAppViewContainer"] > .main {{
                padding-top: 0px;
            }}
            .logo-container {{
                display: flex;
                justify-content: center;
                align-items: center;
                padding: 1rem 2rem;
                padding: 1vh 0; /* margem baseada na altura da tela */
            }}
            .logo-container img {{
                max-height: 20vh; /* altura máxima baseada na tela */
                max-width: 80vw;  /* largura máxima baseada na tela */
                height: auto;
                width: auto;
            }}

            /* Ajuste fino para telas pequenas (MacBook e similares) */
            @media only screen and (max-width: 1440px) {{
                .logo-container {{
                    padding: 3vh 0;
                }}
                .logo-container img {{
                    max-height: 10vh;
                }}
            }}
            </style>
            <div class="logo-container">
                <img src="data:image/png;base64,{encoded_string}" alt="Logo">
            </div>
            """,
            unsafe_allow_html=True
        )